#!/usr/bin/env python3
"""
Build the OSINT for the SE value prop PowerPoint deck.
Uses python-pptx to generate a ~15-slide presentation.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path

# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------

DECK_PATH = Path(__file__).parent / "verkada-se-platform.pptx"

# Colors
NAVY = RGBColor(0x0F, 0x17, 0x2A)
BLUE = RGBColor(0x1E, 0x40, 0xAF)
LIGHT_BLUE = RGBColor(0x60, 0xA5, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_50 = RGBColor(0xF8, 0xFA, 0xFC)
GRAY_100 = RGBColor(0xF1, 0xF5, 0xF9)
GRAY_400 = RGBColor(0x94, 0xA3, 0xB8)
GRAY_600 = RGBColor(0x47, 0x55, 0x63)
GRAY_800 = RGBColor(0x1E, 0x29, 0x3B)
GREEN = RGBColor(0x16, 0xA3, 0x4A)
AMBER = RGBColor(0xD9, 0x77, 0x06)
VIOLET = RGBColor(0x6D, 0x28, 0xD9)


def _add_bg(slide, color=NAVY):
    """Set slide background color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_text_box(slide, left, top, width, height, text, font_size=14,
                  bold=False, color=WHITE, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    """Add a text box with a single run."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox


def _add_multi_text(slide, left, top, width, height, lines, default_size=12, default_color=WHITE):
    """Add a text box with multiple paragraphs (list of dicts or strings)."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        if isinstance(item, str):
            run = p.add_run()
            run.text = item
            run.font.size = Pt(default_size)
            run.font.color.rgb = default_color
            run.font.name = "Calibri"
        elif isinstance(item, dict):
            p.alignment = item.get("align", PP_ALIGN.LEFT)
            p.space_after = Pt(item.get("space_after", 6))
            run = p.add_run()
            run.text = item["text"]
            run.font.size = Pt(item.get("size", default_size))
            run.font.bold = item.get("bold", False)
            run.font.color.rgb = item.get("color", default_color)
            run.font.name = item.get("font", "Calibri")

    return txBox


def _add_rect(slide, left, top, width, height, fill_color, text="",
              font_size=10, font_color=WHITE, bold=False):
    """Add a rounded rectangle with optional text."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False

    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = tf.paragraphs[0].add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = font_color
        run.font.name = "Calibri"

    return shape


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def slide_title(prs):
    """Slide 1: Title"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _add_bg(slide, NAVY)

    _add_text_box(slide, 1.0, 1.5, 8.0, 1.0,
                  "OSINT for the SE", font_size=40, bold=True, color=WHITE)
    _add_text_box(slide, 1.0, 2.5, 8.0, 0.6,
                  "Pre-Call Intelligence for SLED Sales Engineering", font_size=18, color=LIGHT_BLUE)
    _add_text_box(slide, 1.0, 3.5, 8.0, 0.5,
                  "Multi-agent AI research platform  |  18 public data sources  |  1 actionable brief",
                  font_size=12, color=GRAY_400)

    _add_text_box(slide, 1.0, 5.5, 8.0, 0.4,
                  "Andy  |  Solutions Engineer Candidate  |  May 2026",
                  font_size=11, color=GRAY_600)


def slide_problem(prs):
    """Slide 2: The Problem"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, WHITE)

    _add_text_box(slide, 0.5, 0.3, 4.0, 0.4,
                  "THE PROBLEM", font_size=10, bold=True, color=GRAY_400)
    _add_text_box(slide, 0.5, 0.7, 9.0, 0.7,
                  "SE pre-call research is a quota tax", font_size=28, bold=True, color=NAVY)

    _add_multi_text(slide, 0.5, 1.7, 5.5, 3.0, [
        {"text": "Every qualified account requires 2-4 hours of manual research across:", "size": 14, "color": GRAY_800, "space_after": 12},
        {"text": "   SEC filings, news articles, job boards", "size": 12, "color": GRAY_600, "space_after": 4},
        {"text": "   Federal databases (NCES, SAM, HHS, Clery)", "size": 12, "color": GRAY_600, "space_after": 4},
        {"text": "   Cooperative procurement portals", "size": 12, "color": GRAY_600, "space_after": 4},
        {"text": "   Campus safety reports, Reddit, GitHub", "size": 12, "color": GRAY_600, "space_after": 12},
        {"text": "Most SEs skip it. The rest settle for ChatGPT summaries with zero source citations.", "size": 13, "color": GRAY_800, "space_after": 4},
    ])

    # Stats
    _add_rect(slide, 6.5, 1.7, 3.0, 1.2, GRAY_50, bold=True)
    _add_text_box(slide, 6.7, 1.8, 2.6, 0.5, "2-4 hrs", font_size=32, bold=True, color=NAVY)
    _add_text_box(slide, 6.7, 2.35, 2.6, 0.4, "SE prep time per account", font_size=10, color=GRAY_600)

    _add_rect(slide, 6.5, 3.2, 3.0, 1.2, GRAY_50, bold=True)
    _add_text_box(slide, 6.7, 3.3, 2.6, 0.5, "0", font_size=32, bold=True, color=NAVY)
    _add_text_box(slide, 6.7, 3.85, 2.6, 0.4, "source citations from ChatGPT\nor Clay or Sales Navigator", font_size=10, color=GRAY_600)


def slide_why_generic_fails(prs):
    """Slide 3: Why generic AI doesn't solve it"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, WHITE)

    _add_text_box(slide, 0.5, 0.3, 4.0, 0.4,
                  "THE GAP", font_size=10, bold=True, color=GRAY_400)
    _add_text_box(slide, 0.5, 0.7, 9.0, 0.7,
                  "Why generic AI doesn't solve it", font_size=28, bold=True, color=NAVY)

    tools = [
        ("ChatGPT / Claude chat", "No sources. No structure. Could be about any company. Hallucination risk on names, dates, financials."),
        ("Clay / Apollo", "Contact enrichment, not account intelligence. No SLED-specific sources. No synthesis."),
        ("Sales Navigator", "Org chart and news feed. No federal data, procurement signals, or discovery question generation."),
        ("Manual research", "Accurate but takes 2-4 hrs. Doesn't scale. Knowledge walks when the SE does."),
    ]

    y = 1.7
    for tool, gap in tools:
        _add_text_box(slide, 0.8, y, 2.5, 0.4, tool, font_size=13, bold=True, color=NAVY)
        _add_text_box(slide, 3.5, y, 6.0, 0.5, gap, font_size=11, color=GRAY_600)
        y += 0.7


def slide_what_it_is(prs):
    """Slide 4: What this is"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, NAVY)

    _add_text_box(slide, 0.5, 0.3, 4.0, 0.4,
                  "THE SOLUTION", font_size=10, bold=True, color=GRAY_400)
    _add_text_box(slide, 0.5, 0.8, 9.0, 0.8,
                  "OSINT discipline applied to SE pre-call prep",
                  font_size=26, bold=True, color=WHITE)

    _add_multi_text(slide, 0.5, 1.9, 9.0, 3.5, [
        {"text": "Open-source intelligence (OSINT) is the systematic collection and analysis of publicly available information. It's how governments, analysts, and investigators build target profiles.", "size": 13, "color": GRAY_400, "space_after": 16},
        {"text": "This platform applies that discipline to sales engineering:", "size": 14, "color": WHITE, "bold": True, "space_after": 12},
        {"text": "  18 public data sources collected and cached automatically", "size": 12, "color": LIGHT_BLUE, "space_after": 6},
        {"text": "  Multi-agent AI synthesis with anti-genericness rules", "size": 12, "color": LIGHT_BLUE, "space_after": 6},
        {"text": "  Source attribution on every claim (URL + date)", "size": 12, "color": LIGHT_BLUE, "space_after": 6},
        {"text": "  Persona-filtered discovery questions from a rule engine", "size": 12, "color": LIGHT_BLUE, "space_after": 6},
        {"text": "  MEDDIC qualification with named champions", "size": 12, "color": LIGHT_BLUE, "space_after": 6},
        {"text": "  Vendor-agnostic: swap persona file, new vendor", "size": 12, "color": LIGHT_BLUE, "space_after": 6},
    ])


def slide_architecture(prs):
    """Slide 5: Architecture"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, WHITE)

    _add_text_box(slide, 0.5, 0.3, 4.0, 0.4,
                  "ARCHITECTURE", font_size=10, bold=True, color=GRAY_400)
    _add_text_box(slide, 0.5, 0.7, 9.0, 0.5,
                  "Four-phase pipeline", font_size=24, bold=True, color=NAVY)

    phases = [
        ("Phase 1", "COLLECT", "18 source clients\nparallel execution\ncache-aware", NAVY),
        ("Phase 1b", "ENRICH", "Champion signals\nper-individual\nTavily + Haiku", BLUE),
        ("Phase 2", "ANALYZE", "3 Sonnet subagents\nstructured JSON\nconfidence tags", VIOLET),
        ("Phase 3", "SYNTHESIZE", "Opus synthesizer\npersona rules\nMEDDIC + GTM", RGBColor(0x6D, 0x28, 0xD9)),
    ]

    x = 0.5
    for label, title, detail, color in phases:
        _add_rect(slide, x, 1.5, 2.1, 0.4, color, label, font_size=9, font_color=WHITE, bold=True)
        _add_text_box(slide, x, 2.0, 2.1, 0.4, title, font_size=14, bold=True, color=NAVY)
        _add_text_box(slide, x, 2.4, 2.1, 1.0, detail, font_size=10, color=GRAY_600)
        x += 2.3

    # Phase 4
    _add_text_box(slide, 0.5, 3.8, 9.0, 0.4,
                  "Phase 4: RENDER", font_size=14, bold=True, color=NAVY)
    _add_text_box(slide, 0.5, 4.2, 9.0, 0.5,
                  "JSON brief  >  HTML template (Tailwind)  >  Browser-rendered deliverable with source footnote chips",
                  font_size=11, color=GRAY_600)

    # Model assignment
    _add_text_box(slide, 0.5, 5.0, 9.0, 0.3,
                  "MODEL ASSIGNMENT", font_size=9, bold=True, color=GRAY_400)
    _add_text_box(slide, 0.5, 5.3, 9.0, 0.5,
                  "Haiku: parsing/extraction  |  Sonnet: per-source reasoning  |  Opus: final synthesis + specificity rewrite",
                  font_size=10, color=GRAY_600)


def slide_data_layer(prs):
    """Slide 6: The Data Layer"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, WHITE)

    _add_text_box(slide, 0.5, 0.3, 4.0, 0.4,
                  "DATA LAYER", font_size=10, bold=True, color=GRAY_400)
    _add_text_box(slide, 0.5, 0.7, 9.0, 0.5,
                  "18 sources covering SLED-specific intel", font_size=24, bold=True, color=NAVY)

    sources = [
        ("Federal / Compliance", ["SEC EDGAR", "NCES (K-12 districts)", "Clery Act (campus crime)", "SAM.gov (gov registrations)", "HHS OCR (HIPAA breaches)"]),
        ("Market / Procurement", ["Sourcewell coop", "TIPS-USA coop", "OMNIA Partners", "HGACBuy + discounts", "COSTARS (PA)", "GA state procurement", "Atlanta city procurement"]),
        ("Signal / Intelligence", ["News (Tavily)", "Indeed job postings", "crt.sh (SSL/infra)", "GitHub (tech stack)", "Reddit (practitioner)", "Leadership pages"]),
    ]

    x = 0.5
    for category, items in sources:
        _add_text_box(slide, x, 1.4, 2.8, 0.3, category, font_size=11, bold=True, color=BLUE)
        for i, item in enumerate(items):
            _add_text_box(slide, x, 1.8 + i * 0.35, 2.8, 0.3, item, font_size=10, color=GRAY_800)
        x += 3.1


def slide_synthesis(prs):
    """Slide 7: The Synthesis Layer"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, NAVY)

    _add_text_box(slide, 0.5, 0.3, 4.0, 0.4,
                  "SYNTHESIS LAYER", font_size=10, bold=True, color=GRAY_400)
    _add_text_box(slide, 0.5, 0.8, 9.0, 0.5,
                  "The persona file is the leverage", font_size=24, bold=True, color=WHITE)

    rules = [
        "Source attribution required on every factual claim",
        "'Insufficient data' is a valid output — empty beats filler",
        "No generic claims — reject anything that could describe another company",
        "Confidence levels per claim (high / medium / inference)",
        "Trigger-driven discovery questions from persona templates only",
        "Specificity rewrite pass on every synthesis output",
    ]

    for i, rule in enumerate(rules):
        _add_text_box(slide, 0.8, 1.6 + i * 0.55, 8.5, 0.4,
                      f"{i+1}.  {rule}", font_size=12, color=LIGHT_BLUE)

    _add_text_box(slide, 0.5, 5.0, 9.0, 0.5,
                  "verkada-se.yml defines: product lines, ICP verticals, displacement targets,\n"
                  "triggers with discovery templates, buyer personas, disqualifiers, champion criteria",
                  font_size=10, color=GRAY_400)


def slide_demo_setup(prs):
    """Slide 8: Live demo prompt"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, NAVY)

    _add_text_box(slide, 0.5, 0.3, 4.0, 0.4,
                  "LIVE DEMO", font_size=10, bold=True, color=GRAY_400)
    _add_text_box(slide, 0.5, 1.5, 9.0, 0.6,
                  'python3 orchestrator.py "Atlanta Public Schools"',
                  font_size=22, bold=True, color=LIGHT_BLUE)

    _add_text_box(slide, 0.5, 2.5, 9.0, 2.5,
                  "One command.\n\n"
                  "18 source clients fire in parallel.\n"
                  "3 Sonnet subagents analyze in sequence.\n"
                  "Opus synthesizes against the persona rule engine.\n"
                  "HTML brief renders in browser.\n\n"
                  "~3 minutes. Zero manual research.",
                  font_size=14, color=WHITE)


def slide_brief_walkthrough(prs):
    """Slide 9: APS brief highlights"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, WHITE)

    _add_text_box(slide, 0.5, 0.3, 4.0, 0.4,
                  "DEMO: APS BRIEF", font_size=10, bold=True, color=GRAY_400)
    _add_text_box(slide, 0.5, 0.7, 9.0, 0.5,
                  "Three moments to highlight", font_size=24, bold=True, color=NAVY)

    moments = [
        ("1. NDAA + Federal Funding",
         "NCES data: 100% Title I, 73% FRPL, 86 schools. Cross-referenced with Sourcewell #041524-VRK. "
         "Discovery question: 'Has procurement reviewed whether existing security hardware meets NDAA requirements?'"),
        ("2. Cooperative Purchasing Intelligence",
         "HGACBuy: Verkada SE05-26 with Pavion as partner (5% discount tier). Competitor manufacturers visible: "
         "Avigilon, Axis, Genetec, HANWHA, Milestone, Motorola. OMNIA: R250206. This is displacement intel."),
        ("3. Named Champion Projection",
         "Leadership page scraped. Named individuals scored on role fit, tenure, career arc, public voice, "
         "topic affinity, authority. MEDDIC Champion field populated with a person, not a placeholder."),
    ]

    y = 1.5
    for title, detail in moments:
        _add_text_box(slide, 0.8, y, 8.5, 0.3, title, font_size=13, bold=True, color=BLUE)
        _add_text_box(slide, 0.8, y + 0.35, 8.5, 0.7, detail, font_size=10, color=GRAY_600)
        y += 1.3


def slide_meddic(prs):
    """Slide 10: MEDDIC qualification overlay"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, WHITE)

    _add_text_box(slide, 0.5, 0.3, 4.0, 0.4,
                  "MEDDIC", font_size=10, bold=True, color=GRAY_400)
    _add_text_box(slide, 0.5, 0.7, 9.0, 0.5,
                  "Pre-call qualification, not post-call CRM entry", font_size=24, bold=True, color=NAVY)

    fields = [
        ("Metrics", "Quantified from NCES enrollment, site counts, federal funding amounts"),
        ("Economic Buyer", "Named from leadership data + persona meddic_role mappings"),
        ("Decision Criteria", "Derived from vertical drivers + pain hypotheses + NDAA exposure"),
        ("Decision Process", "Inferred from entity type + cooperative purchasing vehicle availability"),
        ("Identify Pain", "Highest-confidence pain hypothesis mapped to Verkada product capability"),
        ("Champion", "Named individual with champion_fit_score, not a role placeholder"),
        ("Competition", "Vendor hits from crt.sh + job postings + cooperative vehicle competitor lists"),
    ]

    y = 1.5
    for field, desc in fields:
        _add_rect(slide, 0.5, y, 1.8, 0.38, NAVY, field, font_size=9, font_color=WHITE, bold=True)
        _add_text_box(slide, 2.5, y, 7.0, 0.4, desc, font_size=10, color=GRAY_600)
        y += 0.55

    _add_text_box(slide, 0.5, 5.5, 9.0, 0.4,
                  "Every field has: evidence array, confidence tag, gap analysis, source quality rating",
                  font_size=10, bold=True, color=GRAY_400)


def slide_champion(prs):
    """Slide 11: Champion projection"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, NAVY)

    _add_text_box(slide, 0.5, 0.3, 4.0, 0.4,
                  "CHAMPION PROJECTION", font_size=10, bold=True, color=GRAY_400)
    _add_text_box(slide, 0.5, 0.8, 9.0, 0.6,
                  "Project, don't detect", font_size=26, bold=True, color=WHITE)

    _add_text_box(slide, 0.5, 1.6, 9.0, 0.5,
                  "Named individuals scored on 6 weighted factors from public signals:",
                  font_size=13, color=GRAY_400)

    factors = [
        ("Role Fit (25%)", "Title parsing against persona role_priority mapping"),
        ("Topic Affinity (20%)", "Haiku topic extraction from public commentary vs champion_topic_signals"),
        ("Career Arc (15%)", "Prior employers cross-referenced against vendor_alumni_indicators"),
        ("Public Voice (15%)", "Speaking, writing, media activity frequency + topics"),
        ("Authority (15%)", "Director+ level, budget mentions, team size hints"),
        ("Recency (10%)", "Tenure inference — recent hires have higher openness signal"),
    ]

    y = 2.3
    for label, desc in factors:
        _add_text_box(slide, 0.8, y, 3.0, 0.3, label, font_size=11, bold=True, color=LIGHT_BLUE)
        _add_text_box(slide, 3.8, y, 5.7, 0.3, desc, font_size=10, color=GRAY_400)
        y += 0.45

    _add_text_box(slide, 0.5, 5.2, 9.0, 0.5,
                  "Output: score_breakdown per factor, reasoning_summary, recommended_validation_question.\n"
                  "These are projections. The SE validates via discovery.",
                  font_size=10, color=GRAY_600)


def slide_gtm(prs):
    """Slide 12: Verkada GTM alignment"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, WHITE)

    _add_text_box(slide, 0.5, 0.3, 4.0, 0.4,
                  "GTM ALIGNMENT", font_size=10, bold=True, color=GRAY_400)
    _add_text_box(slide, 0.5, 0.7, 9.0, 0.5,
                  "Built for Verkada's sales motion", font_size=24, bold=True, color=NAVY)

    items = [
        ("Meraki-style land-and-expand", "Pelican case POC > initial rollout > full platform. Expansion ratio calculated from site count."),
        ("Cooperative-first procurement", "Sourcewell, OMNIA, HGACBuy contracts pre-identified. Bypass full RFP cycle for SLED."),
        ("Displacement intelligence", "crt.sh + job posts + coop vehicle competitor lists surface incumbent vendors. Persona file has verkada_counter text per vendor."),
        ("Channel partner recommendation", "Region + vertical + cooperative data > specific integrator suggestion (e.g., Pavion for GA K-12 via HGACBuy)."),
        ("Bundle recommendation", "Vertical drivers > primary products (Cameras) + expansion products (Access Control, Alarms, Sensors)."),
    ]

    y = 1.5
    for title, desc in items:
        _add_text_box(slide, 0.8, y, 8.5, 0.3, title, font_size=12, bold=True, color=BLUE)
        _add_text_box(slide, 0.8, y + 0.3, 8.5, 0.4, desc, font_size=10, color=GRAY_600)
        y += 0.85


def slide_value_at_scale(prs):
    """Slide 13: Value at scale"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, NAVY)

    _add_text_box(slide, 0.5, 0.3, 4.0, 0.4,
                  "VALUE AT SCALE", font_size=10, bold=True, color=GRAY_400)
    _add_text_box(slide, 0.5, 0.8, 9.0, 0.6,
                  "Time back to quota-carrying activities", font_size=26, bold=True, color=WHITE)

    _add_multi_text(slide, 0.5, 1.8, 9.0, 3.5, [
        {"text": "Per SE: 2-4 hours saved per qualified account", "size": 16, "color": LIGHT_BLUE, "bold": True, "space_after": 12},
        {"text": "10 accounts/month x 3 hrs avg = 30 hrs/month back", "size": 14, "color": WHITE, "space_after": 8},
        {"text": "That's 3.75 selling days recovered per SE per month.", "size": 14, "color": WHITE, "space_after": 20},
        {"text": "At 20 SEs: 750 hours/month = 93 selling days", "size": 14, "color": LIGHT_BLUE, "bold": True, "space_after": 8},
        {"text": "At 100 SEs: 3,750 hours/month = 468 selling days", "size": 14, "color": LIGHT_BLUE, "bold": True, "space_after": 20},
        {"text": "Every hour the SE doesn't spend on research is an hour spent on demos, POCs, and closing.", "size": 12, "color": GRAY_400, "space_after": 4},
    ])


def slide_roadmap(prs):
    """Slide 14: Roadmap"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, WHITE)

    _add_text_box(slide, 0.5, 0.3, 4.0, 0.4,
                  "ROADMAP", font_size=10, bold=True, color=GRAY_400)
    _add_text_box(slide, 0.5, 0.7, 9.0, 0.5,
                  "V1 built. V2 and V3 scoped.", font_size=24, bold=True, color=NAVY)

    tiers = [
        ("V1 — BUILT", GREEN, [
            "18 free public OSINT sources",
            "Multi-agent synthesis (Haiku/Sonnet/Opus)",
            "Persona rule engine (verkada-se.yml)",
            "MEDDIC + champion projection",
            "Cooperative purchasing intelligence",
            "HTML brief with source footnotes",
        ]),
        ("V2 — NEXT", BLUE, [
            "Apollo integration (org charts, direct dials)",
            "Per-district RFP scraping (state procurement portals)",
            "LinkedIn profile enrichment (via API, not scrape)",
            "Automated competitive displacement alerts",
        ]),
        ("V3 — ENTERPRISE", VIOLET, [
            "CRM sync (Salesforce push/pull)",
            "Slack bot for on-demand research",
            "Team dashboards and territory views",
            "Custom vertical modules (healthcare, retail, manufacturing)",
        ]),
    ]

    x = 0.5
    for label, color, items in tiers:
        _add_rect(slide, x, 1.5, 2.8, 0.4, color, label, font_size=10, font_color=WHITE, bold=True)
        for i, item in enumerate(items):
            _add_text_box(slide, x + 0.1, 2.1 + i * 0.4, 2.6, 0.35,
                          f"  {item}", font_size=9, color=GRAY_800)
        x += 3.1


def slide_why_i_built(prs):
    """Slide 15: Why I built this"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, NAVY)

    _add_text_box(slide, 0.5, 0.3, 4.0, 0.4,
                  "WHY I BUILT THIS", font_size=10, bold=True, color=GRAY_400)
    _add_text_box(slide, 0.5, 1.2, 9.0, 0.7,
                  "I wanted to show how I'd actually prepare\nfor accounts if you hire me.",
                  font_size=22, bold=True, color=WHITE)

    _add_multi_text(slide, 0.5, 2.5, 9.0, 3.0, [
        {"text": "This isn't a portfolio project. It's the tool I'd use on day one.", "size": 14, "color": GRAY_400, "space_after": 16},
        {"text": "The persona file is vendor-agnostic by design. verkada-se.yml today, but the architecture works for any SE org that wants sourced, specific, persona-filtered pre-call intelligence.", "size": 13, "color": WHITE, "space_after": 16},
        {"text": "The two demo accounts (Atlanta Public Schools, Georgia Tech) were chosen because they exercise the SLED-specific sources that differentiate this from generic AI research.", "size": 13, "color": WHITE, "space_after": 16},
        {"text": "I built the tool, but the real asset is the methodology.", "size": 14, "color": LIGHT_BLUE, "bold": True, "space_after": 4},
    ])


def slide_qa(prs):
    """Slide 16: Q&A"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, NAVY)

    _add_text_box(slide, 1.0, 2.0, 8.0, 1.0,
                  "Q & A", font_size=48, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    _add_text_box(slide, 1.0, 3.5, 8.0, 1.0,
                  "Pre-loaded for common objections:\n"
                  '"How is this different from Clay?" — Clay is contact enrichment, not account intelligence.\n'
                  '"Does this scale?" — 18 clients run in parallel. ~3 min per account.\n'
                  '"What about data freshness?" — Per-source cache TTLs. News: 7 days. SEC: 90 days.',
                  font_size=11, color=GRAY_400, alignment=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def build():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    slide_title(prs)
    slide_problem(prs)
    slide_why_generic_fails(prs)
    slide_what_it_is(prs)
    slide_architecture(prs)
    slide_data_layer(prs)
    slide_synthesis(prs)
    slide_demo_setup(prs)
    slide_brief_walkthrough(prs)
    slide_meddic(prs)
    slide_champion(prs)
    slide_gtm(prs)
    slide_value_at_scale(prs)
    slide_roadmap(prs)
    slide_why_i_built(prs)
    slide_qa(prs)

    prs.save(str(DECK_PATH))
    print(f"Deck saved: {DECK_PATH} ({DECK_PATH.stat().st_size // 1024}KB)")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    build()
